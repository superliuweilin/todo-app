"""Generate LLM teaching PPT matching the reference '深度学习' PPT style."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy

# ═══════════════════════════════════════════════════════════════
# Color palette (from reference PPT)
# ═══════════════════════════════════════════════════════════════
NAVY        = RGBColor(0x19, 0x2F, 0x59)
LIGHT_BLUE  = RGBColor(0xDB, 0xED, 0xF3)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT  = RGBColor(0xC0, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY    = RGBColor(0x88, 0x88, 0x88)
LIGHTER_BLUE = RGBColor(0xED, 0xF5, 0xFA)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
ORANGE_ACCENT = RGBColor(0xE8, 0x72, 0x2A)
YELLOW_BG   = RGBColor(0xFF, 0xF8, 0xE1)

# 4:3 ratio
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

FONT_CN = "微软雅黑"
FONT_EN = "Times New Roman"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK_LAYOUT = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════
# XML manipulation helpers
# ═══════════════════════════════════════════════════════════════

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

def _make_gradient_xml(top_color, bottom_color, angle=2700000):
    """Create gradient fill XML elements. top_color near stop 0, bottom near stop 1."""
    gsLst = etree.SubElement(etree.Element("dummy"), qn('a:gsLst'))
    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '100000')
    c1 = etree.SubElement(gs1, qn('a:srgbClr'))
    c1.set('val', '{:02X}{:02X}{:02X}'.format(bottom_color[0], bottom_color[1], bottom_color[2]) if isinstance(bottom_color, tuple) else
           '{:02X}{:02X}{:02X}'.format(bottom_color[0] if isinstance(bottom_color, tuple) else 0, 0, 0))

    # Actually let's use RGBColor objects
    return None  # placeholder

def _apply_gradient(shape, color_top, color_bottom, angle=2700000):
    """Apply a linear gradient fill to a shape (via XML manipulation)."""
    nsmap_a = {'a': A_NS}
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = shape._element.find(qn('a:spPr'))
    if spPr is None:
        return

    # Remove existing fill
    for child in list(spPr):
        if child.tag in (qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill')):
            spPr.remove(child)

    # Build gradient fill
    gradFill = etree.SubElement(spPr, qn('a:gradFill'))
    gradFill.set('flip', 'none')
    gradFill.set('rotWithShape', '1')

    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    # Stop 0 (top) - lighter color
    gs0 = etree.SubElement(gsLst, qn('a:gs'))
    gs0.set('pos', '30000')
    c0 = etree.SubElement(gs0, qn('a:srgbClr'))
    c0.set('val', '{:02X}{:02X}{:02X}'.format(color_top[0] if isinstance(color_top, tuple) else color_top[0],
                                                color_top[1] if isinstance(color_top, tuple) else color_top[1],
                                                color_top[2] if isinstance(color_top, tuple) else color_top[2]))

    # Stop 1 (bottom) - darker color
    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '100000')
    c1 = etree.SubElement(gs1, qn('a:srgbClr'))
    c1.set('val', '{:02X}{:02X}{:02X}'.format(color_bottom[0] if isinstance(color_bottom, tuple) else color_bottom[0],
                                                color_bottom[1] if isinstance(color_bottom, tuple) else color_bottom[1],
                                                color_bottom[2] if isinstance(color_bottom, tuple) else color_bottom[2]))

    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '1')
    etree.SubElement(gradFill, qn('a:tileRect'))

def _hex(c):
    """RGBColor -> hex string."""
    if isinstance(c, tuple):
        return '{:02X}{:02X}{:02X}'.format(c[0], c[1], c[2])
    return '{:02X}{:02X}{:02X}'.format(c[0], c[1], c[2])

def _apply_shadow(shape, color=LIGHT_BLUE, alpha=40000, blur=38100):
    """Apply outer shadow effect to a shape."""
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = shape._element.find(qn('a:spPr'))
    if spPr is None:
        return

    # Remove existing effectLst
    for child in list(spPr):
        if child.tag == qn('a:effectLst'):
            spPr.remove(child)

    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('sx', '101000')
    outerShdw.set('sy', '101000')
    outerShdw.set('algn', 'ctr')
    outerShdw.set('rotWithShape', '0')

    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', _hex(color))
    a = etree.SubElement(srgbClr, qn('a:alpha'))
    a.set('val', str(alpha))

def _apply_morph(slide):
    """Apply morph transition to a slide."""
    # Add morph transition
    sld = slide._element
    # Check for existing AlternateContent
    mc_ns = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
    p159_ns = 'http://schemas.microsoft.com/office/powerpoint/2015/09/main'

    nsmap_mc = {'mc': mc_ns}
    ac = etree.SubElement(sld, qn('mc:AlternateContent'))
    choice = etree.SubElement(ac, qn('mc:Choice'))
    choice.set('Requires', 'p159')
    transition = etree.SubElement(choice, qn('p:transition'))
    transition.set('spd', 'slow')
    morph = etree.SubElement(transition, qn('p159:morph'))
    morph.set('option', 'byObject')

    fallback = etree.SubElement(ac, qn('mc:Fallback'))
    trans2 = etree.SubElement(fallback, qn('p:transition'))
    trans2.set('spd', 'slow')
    etree.SubElement(trans2, qn('p:fade'))

# ═══════════════════════════════════════════════════════════════
# High-level shape helpers
# ═══════════════════════════════════════════════════════════════

def add_slide() -> object:
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    _apply_morph(slide)
    return slide

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(14),
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_CN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=Pt(12),
                       color=DARK_GRAY, line_spacing_pts=22):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        if isinstance(line, tuple):
            text, bold, clr = line
            p.text = text
            p.font.bold = bold
            p.font.color.rgb = clr
        else:
            p.text = line
            p.font.color.rgb = color
        p.font.size = font_size
        p.font.name = FONT_CN
    return txBox

def add_page_number(slide, num):
    add_textbox(slide, Inches(9.0), Inches(7.1), Inches(0.7), Inches(0.25),
                str(num), Pt(9), MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle="", page_num=1):
    """Standard section header: navy title + optional subtitle."""
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.5),
                title, Pt(28), NAVY, bold=True)
    # Thin line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.75),
                                   Inches(1.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(9), Inches(0.3),
                    subtitle, Pt(12), MID_GRAY)
    add_page_number(slide, page_num)

def add_gradient_card(slide, left, top, width, height,
                      color_top=WHITE, color_bottom=LIGHT_BLUE):
    """Create a rounded rectangle card with gradient fill + shadow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.line.fill.background()
    _apply_gradient(shape, color_top, color_bottom, 2700000)
    _apply_shadow(shape, LIGHT_BLUE, 40000, 38100)
    return shape

def add_white_label(slide, left, top, width, height, text, font_size=Pt(14)):
    """White rounded rectangle label with shadow (like '知识目标' style)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()
    _apply_shadow(shape, LIGHT_BLUE, 40000, 38100)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_card_with_label(slide, left, top, width, height, label_text, bullets,
                        color_top=WHITE, color_bottom=LIGHT_BLUE, label_width=None):
    """A gradient card with a white label on top, plus bullet content."""
    card = add_gradient_card(slide, left, top, width, height, color_top, color_bottom)
    # Label centered at top of card
    lw = label_width or width - Inches(1.2)
    lh = Inches(0.45)
    lx = left + (width - lw) / 2
    ly = top + Inches(0.1)
    add_white_label(slide, lx, ly, lw, lh, label_text, Pt(13))
    # Bullet content below label
    text_top = ly + lh + Inches(0.12)
    text_height = height - (text_top - top) - Inches(0.1)
    lines = ["• " + b for b in bullets]
    add_multiline_text(slide, left + Inches(0.2), text_top,
                       width - Inches(0.4), text_height,
                       lines, Pt(11), DARK_GRAY, Pt(18))

def add_simple_rect(slide, left, top, width, height, fill_color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_round_rect(slide, left, top, width, height, fill_color, corner=0.08, text="",
                   font_size=Pt(11), font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = corner / 0.166
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_arrow_right(slide, left, top, width, height, fill=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_comparison_table(slide, left, top, col_widths, headers, rows):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    row_h = Inches(0.42)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, row_h * n_rows)
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_CN
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHTER_BLUE if r % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT_CN
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
    return table_shape

# ═══════════════════════════════════════════════════════════════
# Slide builders — 27 slides
# ═══════════════════════════════════════════════════════════════

def s01_cover():
    """Cover slide — dark navy background with white text."""
    slide = add_slide()
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    # Decorative shapes
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(1.0), Inches(3.5), Inches(3.5))
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x22, 0x3F, 0x70); c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), Inches(2.2), Inches(2.0), Inches(2.0))
    c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x2A, 0x4D, 0x80); c2.line.fill.background()
    # Title
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(7), Inches(0.9),
                "\U0001f9e0 大语言模型", Pt(40), WHITE, bold=True)
    add_textbox(slide, Inches(0.6), Inches(2.4), Inches(7), Inches(0.5),
                "Large Language Models", Pt(20), RGBColor(0xAA, 0xBB, 0xDD))
    # Accent line
    add_simple_rect(slide, Inches(0.6), Inches(3.0), Inches(2.0), Inches(0.04), RGBColor(0xAA, 0xBB, 0xDD))
    # Sub info
    add_textbox(slide, Inches(0.6), Inches(3.3), Inches(7), Inches(0.9),
                "\U0001f4da 人工智能概论  ·  专题课程\n\U0001f4a1 从Transformer到GPT-4o：大模型的过去、现在与未来",
                Pt(14), RGBColor(0xCC, 0xDD, 0xEE))

def s02_toc():
    """Table of contents with numbered circles."""
    slide = add_slide()
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(5), Inches(0.55),
                "课程大纲", Pt(28), NAVY, bold=True)
    add_simple_rect(slide, Inches(0.5), Inches(0.7), Inches(1.2), Inches(0.03), NAVY)
    add_page_number(slide, 2)

    sections = [
        ("01", "什么是大语言模型", "定义 · 规模定律 · 涌现能力 · LLM vs 传统NLP"),
        ("02", "核心技术原理", "Transformer · Self-Attention · 训练三阶段 · RLHF"),
        ("03", "主流大模型一览", "GPT系列 · Claude · Gemini · 国产大模型 · 开源生态"),
        ("04", "应用场景", "文本创作 · 编程辅助 · 行业落地 · 多模态 · AI Agent"),
        ("05", "提示工程", "基础原则 · 进阶技巧 · 实践对比 · 好vs坏提示"),
        ("06", "挑战与展望", "幻觉 · 偏见 · 安全 · 能耗 · AGI之路"),
    ]
    for i, (num, title, desc) in enumerate(sections):
        y = Inches(1.2 + i * 1.0)
        # Number circle
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.42), Inches(0.42))
        c.fill.solid(); c.fill.fore_color.rgb = NAVY; c.line.fill.background()
        tf = c.text_frame; p = tf.paragraphs[0]
        p.text = num; p.font.size = Pt(13); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = FONT_CN; p.alignment = PP_ALIGN.CENTER
        add_textbox(slide, Inches(1.2), y + Inches(0.02), Inches(5), Inches(0.28),
                    title, Pt(17), NAVY, bold=True)
        add_textbox(slide, Inches(1.2), y + Inches(0.3), Inches(8), Inches(0.22),
                    desc, Pt(10), MID_GRAY)
        if i < 5:
            add_simple_rect(slide, Inches(1.2), y + Inches(0.62), Inches(8.3), Pt(0.5), LIGHT_BLUE)

def s03_paradigm_shift():
    slide = add_slide()
    add_section_title(slide, "从传统AI到大模型：一次范式革命",
                      "From Narrow AI to Foundation Models", 3)
    eras = [
        ("\U0001f527 专用AI时代\n2010s以前", "每个任务训练一个模型\n需要大量标注数据\n仅限特定领域", LIGHT_BLUE),
        ("\U0001f9e0 深度学习时代\n2012-2020", "CNN/RNN/LSTM\nImageNet推动CV\nBERT/GPT-1初现", LIGHTER_BLUE),
        ("\U0001f680 大模型时代\n2020至今", "一个模型解决多任务\n涌现能力\nGPT-4/Claude/Gemini", LIGHT_BLUE),
    ]
    for i, (title, desc, bg_c) in enumerate(eras):
        left = Inches(0.4 + i * 3.2)
        add_gradient_card(slide, left, Inches(1.3), Inches(2.9), Inches(3.8), WHITE, bg_c)
        add_white_label(slide, left + Inches(0.3), Inches(1.45), Inches(2.3), Inches(0.55),
                        title, Pt(11))
        add_multiline_text(slide, left + Inches(0.25), Inches(2.15), Inches(2.4), Inches(2.8),
                          desc.split("\n"), Pt(12), DARK_GRAY, Pt(20))
        if i < 2:
            add_arrow_right(slide, left + Inches(3.0), Inches(3.0), Inches(0.15), Inches(0.12), NAVY)

def s04_definition():
    slide = add_slide()
    add_section_title(slide, "什么是大语言模型？", "What are Large Language Models?", 4)
    add_card_with_label(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(2.6),
                        "\U0001f4d6 核心概念", [
                            "基于海量文本数据训练的巨型神经网络",
                            "参数规模通常在数十亿到数万亿之间",
                            "核心任务：Next Token Prediction（预测下一个词）",
                            "通过学习语言的统计规律，获得理解和生成能力",
                        ])
    add_card_with_label(slide, Inches(5.2), Inches(1.2), Inches(4.5), Inches(2.6),
                        "⭐ 关键特性", [
                            "规模效应：更大模型 → 更强能力 (Scaling Law)",
                            "上下文学习：通过示例即可学习新任务 (ICL)",
                            "思维链：逐步推理解决复杂问题 (CoT)",
                            "多语言：天然支持多种语言的互译和理解",
                        ])
    # Bottom formula
    fb = add_gradient_card(slide, Inches(0.4), Inches(4.1), Inches(9.3), Inches(1.5), WHITE, LIGHTER_BLUE)
    add_textbox(slide, Inches(0.7), Inches(4.2), Inches(8.8), Inches(0.3),
                "\U0001f4cf 核心公式：P(wₜ | w₁, w₂, ..., wₜ₋₁)", Pt(15), NAVY, bold=True)
    add_multiline_text(slide, Inches(0.7), Inches(4.55), Inches(8.8), Inches(0.9), [
        "大语言模型本质上是一个「条件概率分布函数」—— 给定前文，预测下一个最可能出现的词（token）。",
        "通过反复执行这一简单操作，模型可以写出文章、回答问题、编写代码，甚至进行逻辑推理。",
    ], Pt(11), DARK_GRAY, Pt(18))
    # Flow at bottom
    for i, label in enumerate(["输入文本", "分词Tokenize", "神经网络计算", "输出概率分布", "生成Token"]):
        left = Inches(0.5 + i * 1.95)
        add_round_rect(slide, left, Inches(5.8), Inches(1.65), Inches(0.38), NAVY, text=label, font_size=Pt(10))
        if i < 4:
            add_arrow_right(slide, left + Inches(1.72), Inches(5.9), Inches(0.15), Inches(0.1), MID_GRAY)

def s05_scaling_law():
    slide = add_slide()
    add_section_title(slide, "规模定律：为什么「大」很重要？", "Scaling Law — Kaplan et al., 2020", 5)
    add_card_with_label(slide, Inches(0.4), Inches(1.2), Inches(4.8), Inches(4.5),
                        "\U0001f4c8 Scaling Law 核心发现", [
                            "模型性能与参数量(N)、数据量(D)、计算量(C)呈幂律关系",
                            "性能随规模平滑提升，不存在明显天花板",
                            "更大的模型用更少的数据达到同样性能",
                            "模型大小和数据量应同步扩展",
                            "这一定律推动了模型规模的指数级增长",
                            "从GPT-1的1.17亿→GPT-4的超万亿参数",
                        ])
    # Bar chart on right
    models = [("GPT-1", 0.117), ("BERT-L", 0.34), ("GPT-2", 1.5), ("GPT-3", 175),
              ("PaLM", 540), ("GPT-4", 1800)]
    max_val = 2000
    bar_left = Inches(5.8)
    bar_area_w = Inches(3.6)
    chart_top = Inches(1.4)
    for i, (name, val) in enumerate(models):
        y = chart_top + Inches(i * 0.7)
        bar_w = int(bar_area_w * (val / max_val))
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y, bar_w, Inches(0.42))
        bar.adjustments[0] = 0.3
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        _apply_shadow(bar, NAVY, 20000, 20000)
        add_textbox(slide, bar_left + bar_w + Inches(0.08), y + Inches(0.06),
                    Inches(1.8), Inches(0.3), f"{name}  {val}B", Pt(9), DARK_GRAY, bold=True)

def s06_emergence():
    slide = add_slide()
    add_section_title(slide, "涌现能力：量变引起质变", "Emergent Abilities of LLMs", 6)
    abilities = [
        ("\U0001f9e0", "思维链推理", "模型能分步骤推理\n解决复杂逻辑问题"),
        ("\U0001f4d6", "上下文学习", "通过几个示例\n就能学会新任务"),
        ("\U0001f310", "多语言翻译", "无需平行语料\n实现高质量翻译"),
        ("\U0001f4bb", "代码生成", "理解自然语言\n生成可运行代码"),
        ("\U0001f9ee", "数学推理", "解决需要多步\n计算的数学问题"),
        ("❓", "常识问答", "利用训练习得的\n世界知识回答问题"),
    ]
    for i, (icon, title, desc) in enumerate(abilities):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.2 + row * 3.0)
        add_gradient_card(slide, left, top, Inches(2.9), Inches(2.6), WHITE, LIGHT_BLUE)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.1), Inches(2.6), Inches(0.4),
                    icon, Pt(28))
        add_textbox(slide, left + Inches(0.15), top + Inches(0.65), Inches(2.6), Inches(0.4),
                    title, Pt(15), NAVY, bold=True)
        add_textbox(slide, left + Inches(0.15), top + Inches(1.2), Inches(2.6), Inches(1.2),
                    desc, Pt(11), DARK_GRAY)

def s07_llm_vs_traditional():
    slide = add_slide()
    add_section_title(slide, "大语言模型 vs 传统NLP方法", "LLM vs Traditional NLP — Paradigm Shift", 7)
    col_w = [Inches(2.0), Inches(3.6), Inches(3.6)]
    add_comparison_table(slide, Inches(0.4), Inches(1.4), col_w,
                         ["维度", "传统NLP", "大语言模型"],
                         [
                             ["训练方式", "每个任务单独训练模型", "一个基础模型适配多任务"],
                             ["数据需求", "大量人工标注数据", "海量无标注文本 + 少量标注"],
                             ["泛化能力", "仅限训练过的任务", "可以泛化到未见过的任务"],
                             ["交互方式", "API调用 / 分类输出", "自然语言对话式交互"],
                             ["上下文理解", "固定窗口 / 浅层理解", "长上下文 / 深层语义理解"],
                             ["开发成本", "每个任务高成本", "一次训练，多场景使用"],
                             ["典型代表", "LSTM, CRF, TextCNN", "GPT-4, Claude, Gemini"],
                         ])

def s08_transformer():
    slide = add_slide()
    add_section_title(slide, "Transformer架构：大模型的基石", "Transformer — Attention Is All You Need, 2017", 8)
    # Left: architecture blocks
    blocks = [
        ("输入文本\nToken Embedding", NAVY),
        ("+ 位置编码\nPositional Encoding", RGBColor(0x22, 0x3F, 0x70)),
        ("Multi-Head\nSelf-Attention\n多头自注意力", RGBColor(0x2A, 0x4D, 0x80)),
        ("Feed Forward\n前馈神经网络", NAVY),
        ("输出表示\nOutput", RGBColor(0x22, 0x3F, 0x70)),
    ]
    for i, (text, color) in enumerate(blocks):
        y = Inches(1.3 + i * 0.85)
        add_round_rect(slide, Inches(0.5), y, Inches(2.6), Inches(0.7), color, text=text, font_size=Pt(10))
        if i < 4:
            add_arrow_right(slide, Inches(1.65), y + Inches(0.72), Inches(0.12), Inches(0.1), MID_GRAY)
    # Right: cards
    add_card_with_label(slide, Inches(3.5), Inches(1.3), Inches(6.2), Inches(1.8),
                        "\U0001f4a1 Self-Attention 核心思想", [
                            "让每个词都关注到序列中所有其他词，计算它们之间的关联权重",
                            "Q(Query) × K(Key)ᵀ → Softmax → × V(Value) → 加权求和",
                            "相比RNN：可并行计算、能捕捉长距离依赖、训练效率高",
                        ])
    add_card_with_label(slide, Inches(3.5), Inches(3.4), Inches(6.2), Inches(1.8),
                        "\U0001f527 Transformer 关键创新", [
                            "抛弃RNN/LSTM的序列计算方式，完全基于注意力机制",
                            "Encoder（编码器）：理解输入 → Decoder（解码器）：生成输出",
                            "现代LLM大多使用Decoder-Only架构（如GPT系列）",
                        ])
    # Bottom timeline
    steps = ["Transformer\n2017", "BERT\n2018", "GPT-2\n2019", "GPT-3\n2020", "ChatGPT\n2022", "GPT-4\n2023", "GPT-4o\n2024"]
    for i, s in enumerate(steps):
        left = Inches(0.4 + i * 1.38)
        add_round_rect(slide, left, Inches(5.5), Inches(1.15), Inches(0.55), NAVY, text=s, font_size=Pt(8))

def s09_attention():
    slide = add_slide()
    add_section_title(slide, "Self-Attention 机制详解", "How Self-Attention Works — Step by Step", 9)
    add_multiline_text(slide, Inches(0.4), Inches(1.1), Inches(9.5), Inches(0.8), [
        ("Self-Attention 计算步骤：", True, NAVY),
        "① 每个词生成 Q(Query)、K(Key)、V(Value) 三个向量    ② 用Q和K计算注意力分数矩阵",
        "③ Softmax归一化得到注意力权重    ④ 用权重对V加权求和，得到每个词的上下文表示",
    ], Pt(11), DARK_GRAY, Pt(18))
    # Example
    words = ["中国", "的", "首都", "是", "北京"]
    n = len(words)
    cell_w = Inches(1.6)
    start_left = Inches(0.8)
    for i, w in enumerate(words):
        left = start_left + cell_w * i
        add_round_rect(slide, left, Inches(2.4), Inches(1.4), Inches(0.45), NAVY, text=f"Q: {w}", font_size=Pt(11))
    # Matrix
    for row in range(n):
        for col in range(n):
            left = start_left + cell_w * col
            top = Inches(3.1 + row * 0.35)
            if row == col:
                c = RGBColor(0x30, 0x60, 0xA0)
            elif abs(row - col) <= 1:
                c = RGBColor(0x70, 0xA0, 0xD0)
            else:
                c = LIGHT_BLUE
            add_simple_rect(slide, left, top, Inches(1.4), Inches(0.3), c)
    # Output
    for i, w in enumerate(words):
        left = start_left + cell_w * i
        add_round_rect(slide, left, Inches(5.1), Inches(1.4), Inches(0.45), RGBColor(0x2A, 0x4D, 0x80), text=f"Out: {w}", font_size=Pt(11))
    add_textbox(slide, Inches(0.4), Inches(5.8), Inches(9.5), Inches(0.3),
                "说明：每个词关注所有其他词，但注意力权重不同。'北京'和'首都'之间权重较高，因为它们语义相关。",
                Pt(10), MID_GRAY)

def s10_training():
    slide = add_slide()
    add_section_title(slide, "大模型的训练三阶段", "Training Pipeline: Pre-training → SFT → RLHF", 10)
    stages = [
        ("\U0001f4da 阶段一：预训练\nPre-training",
         ["目标：学习语言的统计规律", "数据：数万亿Token的文本", "方法：自监督 Next Token Prediction",
          "算力：数千GPU，数月训练", "产出：基座模型 (Base Model)", "代表：GPT-3, LLaMA, Qwen"]),
        ("\U0001f3af 阶段二：监督微调\nSFT",
         ["目标：学会遵循指令和对话", "数据：高质量问答对 (10万-100万)", "方法：标准监督学习",
          "产出：对话模型 (Chat Model)", "代表：GPT-3.5-turbo, Vicuna"]),
        ("⚖️ 阶段三：RLHF\n人类反馈强化学习",
         ["目标：对齐人类偏好和价值观", "数据：人类偏好比较数据", "方法：PPO / DPO 算法",
          "关键：奖励模型 + 策略优化", "产出：对齐模型 (Aligned Model)", "代表：ChatGPT, Claude"]),
    ]
    for i, (title, bullets) in enumerate(stages):
        left = Inches(0.3 + i * 3.25)
        add_card_with_label(slide, left, Inches(1.3), Inches(3.0), Inches(4.5), title, bullets,
                           color_bottom=LIGHT_BLUE if i != 2 else RGBColor(0xFF, 0xEE, 0xDD))
        if i < 2:
            add_arrow_right(slide, left + Inches(3.1), Inches(3.2), Inches(0.13), Inches(0.1), NAVY)

def s11_rlhf():
    slide = add_slide()
    add_section_title(slide, "RLHF：让模型对齐人类偏好", "Reinforcement Learning from Human Feedback", 11)
    boxes = [("预训练\n基座模型", NAVY), ("收集人类\n偏好数据", RGBColor(0x22, 0x3F, 0x70)),
             ("训练\n奖励模型", RGBColor(0x2A, 0x4D, 0x80)), ("PPO/DPO\n策略优化", RGBColor(0x30, 0x55, 0x90)),
             ("对齐后\n的模型", RGBColor(0x2C, 0x6C, 0x4C))]
    for i, (text, color) in enumerate(boxes):
        left = Inches(0.3 + i * 1.95)
        add_round_rect(slide, left, Inches(1.4), Inches(1.6), Inches(0.9), color, text=text, font_size=Pt(11))
        if i < 4:
            add_arrow_right(slide, left + Inches(1.72), Inches(1.7), Inches(0.15), Inches(0.1), MID_GRAY)

    add_card_with_label(slide, Inches(0.4), Inches(2.7), Inches(4.5), Inches(3.2),
                        "\U0001f504 RLHF 核心流程", [
                            "1. 从预训练SFT模型采样多个回答",
                            "2. 人类标注者对回答进行排序（A > B > C）",
                            "3. 训练奖励模型(RM)预测人类偏好分数",
                            "4. 使用PPO算法优化策略，最大化RM打分",
                            "5. 加入KL散度约束，防止偏离太远",
                            "6. DPO简化了PPO，直接优化偏好数据",
                        ])
    add_card_with_label(slide, Inches(5.2), Inches(2.7), Inches(4.5), Inches(3.2),
                        "❓ 为什么需要RLHF？", [
                            "基座模型会生成有害/偏见内容",
                            "模型不知道什么是「好」的回答",
                            "RLHF让模型学习人类价值观",
                            "提升有用性(Helpful)和安全性(Harmless)",
                            "减少幻觉和拒绝不当请求",
                            "这也被称为 AI对齐 (Alignment)",
                        ])

def s12_tokenization():
    slide = add_slide()
    add_section_title(slide, "Token与分词：模型如何「读」文字", "Tokenization — The First Step of LLM", 12)
    # Example
    eg = add_gradient_card(slide, Inches(0.4), Inches(1.2), Inches(9.3), Inches(0.7), WHITE, LIGHTER_BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(8.8), Inches(0.5),
                '\U0001f4dd 原文："人工智能正在改变世界"  →  Token序列：[人工智能, 正在, 改变, 世界]',
                Pt(15), DARK_GRAY, bold=True)
    # Comparison table
    col_w = [Inches(1.8), Inches(2.5), Inches(2.5), Inches(2.5)]
    add_comparison_table(slide, Inches(0.4), Inches(2.2), col_w,
                         ["", "BPE (GPT系列)", "SentencePiece (LLaMA)", "WordPiece (BERT)"],
                         [
                             ["原理", "字节对编码\n高频字符合并", "Unigram语言模型\n概率分词", "贪心最长匹配\n字符级BPE"],
                             ["词汇量", "~100K", "~32K-256K", "~30K"],
                             ["中文处理", "按UTF-8字节\n效果一般", "支持中文\n效果较好", "按字符分割\n效率较低"],
                         ])
    add_card_with_label(slide, Inches(0.4), Inches(4.1), Inches(9.3), Inches(1.7),
                        "⚠️ 为什么Tokenization很重要？", [
                            "Token决定了模型理解的「最小单元」——类似于人类语言中的「词」或「字」",
                            "同样的中文输入，不同分词器产生的Token数量不同 → 影响推理速度和成本",
                            "GPT-4的中文Token效率约是英文的2倍（同样内容中文消耗更多Token）",
                            "分词器的不完美会导致「Token Bias」——例如数字分解后算术推理变得困难",
                        ])

def s13_gpt_series():
    slide = add_slide()
    add_section_title(slide, "GPT系列：从实验室到席卷全球", "The GPT Family — A 6-Year Journey", 13)
    models = [
        ("GPT-1\n2018", "1.17亿参数\nBooksCorpus\n概念验证期"),
        ("GPT-2\n2019", "15亿参数\n\"Too dangerous\"\n引发AI安全讨论"),
        ("GPT-3\n2020", "1750亿参数\nIn-Context Learning\nFew-Shot能力"),
        ("GPT-3.5\n2022", "ChatGPT发布\nRLHF对齐\n月活1亿/2个月"),
        ("GPT-4\n2023", "多模态(MoE)\n律师考试前10%\n更大上下文"),
        ("GPT-4o\n2024", "全能多模态\n实时语音/视觉\n更快更便宜"),
    ]
    for i, (name, desc) in enumerate(models):
        left = Inches(0.25 + i * 1.62)
        add_round_rect(slide, left, Inches(1.3), Inches(1.35), Inches(0.8), NAVY, text=name, font_size=Pt(10))
        add_gradient_card(slide, left, Inches(2.3), Inches(1.35), Inches(2.1), WHITE, LIGHT_BLUE)
        add_multiline_text(slide, left + Inches(0.08), Inches(2.4), Inches(1.2), Inches(1.9),
                          desc.split("\n"), Pt(9), DARK_GRAY, Pt(14))
        if i < 5:
            add_arrow_right(slide, left + Inches(1.42), Inches(1.65), Inches(0.12), Inches(0.08), MID_GRAY)
    # Bottom
    add_textbox(slide, Inches(0.4), Inches(4.7), Inches(9.5), Inches(0.3),
                "\U0001f4a1 关键启示：", Pt(14), NAVY, bold=True)
    add_multiline_text(slide, Inches(0.4), Inches(5.0), Inches(9.3), Inches(2.0), [
        "从GPT-1到GPT-4o，仅6年时间，参数规模增长约15000倍",
        "GPT-3验证了Scaling Law的正确性，Few-Shot能力改变了对AI的认知",
        "ChatGPT将大模型推向大众，2个月内月活突破1亿，成为史上增长最快的应用",
        "GPT-4引入多模态能力（理解图像），GPT-4o进一步统一文本/视觉/语音",
        "MoE（混合专家）架构：GPT-4据传由8个×220B专家子模型组成，推理时仅激活部分参数",
    ], Pt(10), DARK_GRAY, Pt(16))

def s14_claude_gemini():
    slide = add_slide()
    add_section_title(slide, "Claude 与 Gemini：两大强力竞争者", "Major Competitors: Anthropic Claude & Google Gemini", 14)
    add_card_with_label(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(2.6),
                        "\U0001f916 Anthropic — Claude 系列", [
                            "由前OpenAI员工创立，专注AI安全",
                            "Claude 3.5 Sonnet（2024.6）：编程能力顶级",
                            "Claude 3.5 Haiku：快速且经济",
                            "Claude Opus 4（2025）：旗舰级推理能力",
                            "核心特点：长上下文(200K)、安全性强",
                            "Constitutional AI 对齐方法",
                        ])
    add_card_with_label(slide, Inches(5.2), Inches(1.2), Inches(4.5), Inches(2.6),
                        "\U0001f310 Google — Gemini 系列", [
                            "DeepMind + Google Brain 合并打造",
                            "Gemini 1.0（2023.12）：原生多模态",
                            "Gemini 1.5 Pro：100万Token上下文窗口",
                            "Gemini 2.5 Pro（2025）：推理能力领先",
                            "核心特点：多模态原生融合、超长上下文",
                            "深度整合Google生态（搜索、Workspace）",
                        ])
    col_w = [Inches(2.2), Inches(2.5), Inches(2.4), Inches(2.4)]
    add_comparison_table(slide, Inches(0.4), Inches(4.2), col_w,
                         ["维度", "OpenAI (GPT-4o)", "Anthropic (Claude)", "Google (Gemini)"],
                         [
                             ["核心优势", "多模态+生态\n开发者工具多", "长上下文+安全\n编程/写作强", "超长上下文\n搜索整合"],
                             ["上下文窗口", "128K", "200K", "1M tokens"],
                             ["安全理念", "RLHF+红队测试", "Constitutional AI", "AI Principles"],
                             ["定价策略", "中高", "中等", "中等偏低"],
                         ])

def s15_chinese_models():
    slide = add_slide()
    add_section_title(slide, "国产大模型：百花齐放", "Chinese LLMs Landscape \U0001f1e8\U0001f1f3", 15)
    models = [
        ("百度 · 文心一言", "ERNIE系列 | 知识增强\n中文理解领先\n2023.3首发"),
        ("阿里 · 通义千问", "Qwen系列 | 开源生态好\n多模态 | 中文能力强\n全球下载量领先"),
        ("DeepSeek", "MoE架构 | 高性价比\nDeepSeek-V3/R1\n推理能力极强"),
        ("字节 · 豆包", "火山引擎驱动\n模型家族完整\n语音助手领先"),
        ("月之暗面 · Kimi", "超长上下文200万字\n联网搜索\n深度阅读"),
        ("智谱 · ChatGLM", "清华系 | GLM架构\n中英双语\n完整产品线"),
    ]
    for i, (name, desc) in enumerate(models):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.2 + row * 3.0)
        add_card_with_label(slide, left, top, Inches(2.9), Inches(2.6), name, desc.split("\n"),
                           color_bottom=LIGHT_BLUE if i != 2 else RGBColor(0xFF, 0xEE, 0xDD))

def s16_open_source():
    slide = add_slide()
    add_section_title(slide, "开源大模型生态", "Open Source LLM Ecosystem", 16)
    col_w = [Inches(1.6), Inches(1.8), Inches(1.5), Inches(2.0), Inches(2.5)]
    add_comparison_table(slide, Inches(0.25), Inches(1.3), col_w,
                         ["模型", "开发者", "参数规模", "核心特色", "适用场景"],
                         [
                             ["LLaMA 3/4", "Meta", "8B-405B", "开源标杆，生态最完善", "研究、微调、多场景"],
                             ["Qwen 2.5", "阿里", "0.5B-72B", "中文最强，多尺寸覆盖", "中文应用、企业部署"],
                             ["Mistral", "Mistral AI", "7B-8×22B", "MoE架构，法文优化", "欧洲语言、轻量场景"],
                             ["DeepSeek-V3", "DeepSeek", "671B(MoE)", "创新MoE，高性价比", "中文推理、编程"],
                             ["Phi-4", "微软", "14B", "小模型大能力", "端侧部署、教育场景"],
                         ])
    add_card_with_label(slide, Inches(0.4), Inches(4.2), Inches(9.3), Inches(2.0),
                        "\U0001f4a1 开源模型的意义与趋势", [
                            "降低AI门槛：开发者可免费使用、微调、部署，无需依赖API，隐私性更好",
                            "缩小差距：开源模型与闭源模型的性能差距正在快速缩小（Qwen2.5-72B ≈ GPT-4约85%）",
                            "垂直应用：企业可基于开源模型微调行业专用模型（法律、医疗、金融等）",
                            "小型化趋势：1B-7B参数的「小钢炮」模型能力持续提升，端侧部署成为可能",
                        ])

def s17_comparison():
    slide = add_slide()
    add_section_title(slide, "主流大模型能力全景对比", "Model Capability Comparison — 2025.5", 17)
    col_w = [Inches(1.7), Inches(1.3), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.1), Inches(1.1)]
    add_comparison_table(slide, Inches(0.3), Inches(1.3), col_w,
                         ["模型", "参数", "上下文", "多模态", "编程", "中文", "价格"],
                         [
                             ["GPT-4o", "~1.8T(MoE)", "128K", "文+图+音", "★★★★★", "★★★★", "$$$"],
                             ["GPT-4.1", "~1.8T(MoE)", "1M", "文+图", "★★★★★", "★★★★", "$$$"],
                             ["Claude Opus 4", "未公开", "200K", "文+图", "★★★★★", "★★★★", "$$$"],
                             ["Claude Sonnet 4", "未公开", "200K", "文+图", "★★★★☆", "★★★★", "$$"],
                             ["Gemini 2.5 Pro", "未公开", "1M", "原生多模态", "★★★★☆", "★★★☆", "$$"],
                             ["DeepSeek-V3", "671B(MoE)", "128K", "文本", "★★★★★", "★★★★★", "$"],
                             ["DeepSeek-R1", "671B(MoE)", "128K", "文本", "★★★★☆", "★★★★★", "$"],
                             ["Qwen2.5-72B", "72B", "128K", "文+图", "★★★★", "★★★★★", "$"],
                         ])
    add_textbox(slide, Inches(0.4), Inches(6.8), Inches(9.5), Inches(0.3),
                "注：价格相对参考 ($低  $$$高)  |  数据截至2025年5月  |  开源模型标注下划线",
                Pt(9), MID_GRAY)

def s18_text_gen():
    slide = add_slide()
    add_section_title(slide, "应用场景（一）：文本生成与创作", "Applications: Text Generation & Creation", 18)
    apps = [
        ("\U0001f4dd 内容创作", ["文章写作、剧本创作", "诗歌、小说、广告文案", "多语言翻译与润色"]),
        ("\U0001f4ca 信息处理", ["文档摘要与提炼", "会议纪要自动生成", "长篇报告分析"]),
        ("\U0001f4ac 对话与客服", ["智能客服 7×24小时", "个性化教育辅导", "心理咨询辅助"]),
        ("\U0001f4e7 办公效率", ["邮件起草与回复", "PPT大纲生成", "Excel公式编写"]),
    ]
    for i, (title, bullets) in enumerate(apps):
        col, row = i % 2, i // 2
        left = Inches(0.4 + col * 4.9)
        top = Inches(1.2 + row * 2.9)
        add_card_with_label(slide, left, top, Inches(4.5), Inches(2.5), title, bullets)

def s19_coding():
    slide = add_slide()
    add_section_title(slide, "应用场景（二）：编程与代码辅助", "Applications: Coding & Development", 19)
    add_card_with_label(slide, Inches(0.4), Inches(1.2), Inches(2.9), Inches(3.5),
                        "\U0001f916 GitHub Copilot", [
                            "AI编程助手先驱",
                            "代码补全与生成",
                            "支持VS Code/JetBrains",
                            "自然语言→代码",
                            "覆盖数十种编程语言",
                            "用户超200万(2024)",
                        ])
    add_card_with_label(slide, Inches(3.55), Inches(1.2), Inches(2.9), Inches(3.5),
                        "\U0001f4bb Cursor / Windsurf", [
                            "AI原生IDE",
                            "整文件/多文件编辑",
                            "自然语言对话式编程",
                            "上下文感知整个项目",
                            "一键修复Bug",
                            "新一代编程范式",
                        ])
    add_card_with_label(slide, Inches(6.7), Inches(1.2), Inches(2.9), Inches(3.5),
                        "☁️ 通用模型编程", [
                            "Claude：编程能力最强",
                            "GPT-4o：代码解释完善",
                            "DeepSeek-V3：高性价比",
                            "通过对话开发应用",
                            "原型开发效率提升10x",
                            "非程序员也能写代码",
                        ])
    add_textbox(slide, Inches(0.4), Inches(4.9), Inches(9.5), Inches(0.3),
                "\U0001f504 典型工作流：", Pt(14), NAVY, bold=True)
    steps = ["描述需求", "AI生成代码", "审查测试", "迭代修改", "完成部署"]
    for i, s in enumerate(steps):
        left = Inches(0.4 + i * 1.95)
        add_round_rect(slide, left, Inches(5.3), Inches(1.65), Inches(0.42), NAVY, text=s, font_size=Pt(10))
        if i < 4:
            add_arrow_right(slide, left + Inches(1.72), Inches(5.42), Inches(0.15), Inches(0.1), MID_GRAY)

def s20_industries():
    slide = add_slide()
    add_section_title(slide, "应用场景（三）：行业落地", "Applications: Industry Verticals", 20)
    industries = [
        ("\U0001f3e5 医疗", ["辅助诊断", "病历分析", "药物研发", "医学文献检索"]),
        ("⚖️ 法律", ["合同审查", "案例检索", "法律文书", "合规检查"]),
        ("\U0001f3e6 金融", ["风险评估", "投资分析", "反欺诈", "智能投顾"]),
        ("\U0001f393 教育", ["个性化辅导", "自动出题", "作文批改", "知识问答"]),
        ("\U0001f52c 科研", ["文献综述", "实验设计", "数据分析", "论文润色"]),
        ("\U0001f3ed 制造", ["质量检测", "故障预测", "供应链优化", "操作手册"]),
    ]
    for i, (title, bullets) in enumerate(industries):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.2 + row * 3.0)
        add_card_with_label(slide, left, top, Inches(2.9), Inches(2.6), title, bullets)

def s21_multimodal():
    slide = add_slide()
    add_section_title(slide, "应用场景（四）：多模态能力", "Applications: Multimodal Capabilities", 21)
    mods = [
        ("\U0001f5bc️ 图像理解", ["识别图片内容、OCR文字提取", "图表分析、医学影像辅助", "GPT-4V/4o, Claude, Gemini 支持"]),
        ("\U0001f3a4 语音交互", ["语音→文字（ASR）", "文字→语音（TTS）", "实时语音对话：GPT-4o 高级语音模式"]),
        ("\U0001f3ac 视频理解", ["视频内容摘要", "动作识别与事件检测", "Gemini 1.5 Pro 领先"]),
        ("\U0001f3a8 图像生成", ["文字→图像（DALL-E, Midjourney）", "风格迁移、图像编辑", "GPT-4o 内置图像生成"]),
    ]
    for i, (title, bullets) in enumerate(mods):
        col, row = i % 2, i // 2
        left = Inches(0.4 + col * 4.9)
        top = Inches(1.2 + row * 2.9)
        add_card_with_label(slide, left, top, Inches(4.5), Inches(2.5), title, bullets)

def s22_agents():
    slide = add_slide()
    add_section_title(slide, "应用场景（五）：AI Agent 智能体", "Applications: AI Agents — Autonomous AI", 22)
    add_card_with_label(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(3.0),
                        "\U0001f916 什么是AI Agent？", [
                            "AI Agent是能自主感知环境、制定计划、",
                            "使用工具、执行行动来完成复杂任务的AI系统",
                            "",
                            "核心理念：LLM作为「大脑」+ 工具调用能力",
                            "→ 从被动问答到主动执行任务",
                            "",
                            "范式升级：Copilot（副驾驶）→ Agent（自动驾驶）",
                        ])
    add_card_with_label(slide, Inches(5.2), Inches(1.2), Inches(4.5), Inches(3.0),
                        "⚙️ Agent的核心能力", [
                            "规划 (Planning)：任务分解与步骤规划",
                            "记忆 (Memory)：短期/长期记忆管理",
                            "工具使用 (Tool Use)：API、搜索、代码执行",
                            "反思 (Reflection)：自我检查与纠错",
                            "多Agent协作：多个Agent分工合作",
                        ])
    # Flow
    add_textbox(slide, Inches(0.4), Inches(4.5), Inches(9.5), Inches(0.3),
                "\U0001f504 Agent工作流程：", Pt(13), NAVY, bold=True)
    flow = ["用户目标", "任务规划", "调用工具", "执行行动", "观察结果", "反思调整", "完成目标"]
    for i, f in enumerate(flow):
        left = Inches(0.3 + i * 1.4)
        add_round_rect(slide, left, Inches(4.9), Inches(1.15), Inches(0.5), NAVY if i < 4 else RGBColor(0x2C, 0x6C, 0x4C),
                       text=f, font_size=Pt(9))
        if i < 6:
            add_arrow_right(slide, left + Inches(1.22), Inches(5.05), Inches(0.12), Inches(0.08), MID_GRAY)

def s23_prompt_basics():
    slide = add_slide()
    add_section_title(slide, "提示工程（一）：基础原则", "Prompt Engineering — Basic Principles", 23)
    add_card_with_label(slide, Inches(0.4), Inches(1.2), Inches(3.0), Inches(3.2),
                        "\U0001f4a1 原则一：清晰明确", [
                            "明确指定任务和目标",
                            "给出具体的输出格式",
                            "设定角色和背景",
                            "\"写一份报告\" → ✗",
                            "\"以数据分析师身份，用",
                            "Markdown格式写一份2024",
                            "AI行业趋势分析报告\" → ✓",
                        ])
    add_card_with_label(slide, Inches(3.65), Inches(1.2), Inches(3.0), Inches(3.2),
                        "\U0001f4dd 原则二：提供上下文", [
                            "说明你的背景和需求",
                            "给出相关领域信息",
                            "指定受众和语言风格",
                            "\"我是大一学生，请用通俗",
                            "易懂的方式解释Transformer",
                            "架构，可以打比方\"",
                        ])
    add_card_with_label(slide, Inches(6.9), Inches(1.2), Inches(2.9), Inches(3.2),
                        "\U0001f504 原则三：迭代优化", [
                            "先给简单版本试效果",
                            "根据输出逐步细化",
                            "把复杂任务拆解成多步",
                            "使用约束条件控制输出",
                            "\"请用JSON格式输出，",
                            "包含字段：name, summary,",
                            "key_points（3-5条）\"",
                        ])
    # Formula
    add_textbox(slide, Inches(0.4), Inches(4.7), Inches(9.5), Inches(0.3),
                "\U0001f4cf 提示词 = 指令 + 上下文 + 输入数据 + 输出格式要求", Pt(14), NAVY, bold=True)
    formula = [("指令\nInstruction", NAVY), ("上下文\nContext", RGBColor(0x22, 0x3F, 0x70)),
               ("输入\nInput", RGBColor(0x2A, 0x4D, 0x80)), ("格式\nFormat", RGBColor(0x30, 0x55, 0x90)),
               ("有效提示词\nPrompt", RGBColor(0x2C, 0x6C, 0x4C))]
    for i, (text, color) in enumerate(formula):
        left = Inches(0.4 + i * 1.95)
        add_round_rect(slide, left, Inches(5.15), Inches(1.65), Inches(0.55), color, text=text, font_size=Pt(9))

def s24_prompt_advanced():
    slide = add_slide()
    add_section_title(slide, "提示工程（二）：进阶技巧", "Prompt Engineering — Advanced Techniques", 24)
    techniques = [
        ("\U0001f9e9 思维链(CoT)", ["让模型「一步步思考」", "\"Let's think step by step\"", "显著提升推理准确性", "数学题、逻辑题必备"]),
        ("\U0001f4cb 少样本(Few-shot)", ["在Prompt中给2-5个示例", "模型模仿示例格式和风格", "无需微调即可定制输出", "最实用的技巧之一"]),
        ("\U0001f4e6 结构化输出", ["指定JSON/XML/Markdown格式", "约束输出字段和类型", "方便程序解析和自动化", "生产环境必备"]),
        ("\U0001f333 思维树(ToT)", ["探索多条推理路径", "评估并选择最佳路径", "回溯和修正能力", "适合复杂决策问题"]),
        ("\U0001f50d RAG", ["外挂知识库/文档", "先检索相关信息再回答", "减少幻觉、实时更新", "企业应用主流方案"]),
        ("\U0001f916 ReAct", ["推理+行动交替", "Thought→Action→Observation", "适合Agent场景", "工具调用+反思循环"]),
    ]
    for i, (title, bullets) in enumerate(techniques):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.2 + row * 3.0)
        add_card_with_label(slide, left, top, Inches(2.9), Inches(2.6), title, bullets,
                           color_bottom=LIGHT_BLUE if i != 4 else YELLOW_BG)

def s25_prompt_examples():
    slide = add_slide()
    add_section_title(slide, "提示工程（三）：实践对比", "Prompt Examples — Good vs Bad", 25)
    # Bad
    bad_card = add_gradient_card(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(2.6),
                                 WHITE, RGBColor(0xFD, 0xED, 0xED))
    add_white_label(slide, Inches(0.7), Inches(1.35), Inches(3.9), Inches(0.42),
                    "❌  不好的提示", Pt(13))
    add_multiline_text(slide, Inches(0.7), Inches(1.95), Inches(3.9), Inches(1.7), [
        "\"帮我写个关于气候变化的文章\"",
        "",
        "问题：",
        "没有指定长度、风格、受众",
        "没有说明具体角度或重点",
        "输出难以控制，可能不符合需求",
    ], Pt(10), DARK_GRAY, Pt(15))
    # Good
    good_card = add_gradient_card(slide, Inches(5.2), Inches(1.2), Inches(4.5), Inches(2.6),
                                  WHITE, RGBColor(0xD4, 0xED, 0xDA))
    add_white_label(slide, Inches(5.5), Inches(1.35), Inches(3.9), Inches(0.42),
                    "✅  好的提示", Pt(13))
    add_multiline_text(slide, Inches(5.5), Inches(1.95), Inches(3.9), Inches(1.7), [
        "\"你是一位环境科学教授，请为大一学生写",
        "一篇800字的科普短文，介绍气候变化",
        "对中国农业的影响。要求：语言通俗易懂，",
        "包含≥3个具体数据，用小标题分段，",
        "结尾提出2-3个可行的个人行动建议。\"",
        "",
        "优点：角色设定、受众、长度、内容要素、",
        "结构要求都明确，一次生成即可满足需求",
    ], Pt(10), DARK_GRAY, Pt(15))
    # Bottom
    add_card_with_label(slide, Inches(0.4), Inches(4.1), Inches(9.3), Inches(2.2),
                        "\U0001f4a1 提示工程实践建议", [
                            "角色扮演往往效果显著：\"你是一位资深Python工程师\"比直接提问获得更好的代码",
                            "让模型自我修正：在提示结尾加 \"请检查你的回答是否有误，如有错误请自行修正\"",
                            "分步复杂任务：\"首先列出大纲 → 我确认 → 再逐节展开\" 比一次性生成更有质量保证",
                            "善用约束条件：字数/格式/风格/禁用词/必须包含的要点 —— 越具体越好",
                            "提示词也需要版本管理！好的提示词值得保存和迭代，这是与大模型协作的核心技能",
                        ])

def s26_challenges():
    slide = add_slide()
    add_section_title(slide, "挑战与风险：大模型并非完美", "Challenges & Risks — LLMs Are Not Perfect", 26)
    risks = [
        ("\U0001f534 幻觉 Hallucination", "模型会自信地生成不存在的事实\n学术写作需要核实 | RAG可部分缓解"),
        ("\U0001f7e0 偏见 Bias & Fairness", "训练数据中的偏见被模型放大\n性别/种族/地域偏见 | 去偏技术探索中"),
        ("\U0001f7e1 安全 Safety & Misuse", "虚假信息/钓鱼邮件/深度伪造\n越狱攻击(Jailbreak) | 需安全防护"),
        ("\U0001f7e2 隐私 Privacy & IP", "训练数据版权争议\n用户隐私泄露 | 法律框架不完善"),
        ("\U0001f535 能耗 Energy & Cost", "训练一次耗电巨大\n推理需大量算力 | 小模型/蒸馏是方向"),
        ("⚫ 不可控性 Uncertainty", "输出难以精确控制\n复杂推理易出错 | 可解释性不足"),
    ]
    for i, (title, desc) in enumerate(risks):
        col, row = i % 3, i // 3
        left = Inches(0.4 + col * 3.2)
        top = Inches(1.2 + row * 3.0)
        add_card_with_label(slide, left, top, Inches(2.9), Inches(2.6), title, desc.split("\n"))

def s27_future():
    slide = add_slide()
    add_section_title(slide, "未来展望：AGI还有多远？", "The Road to AGI — Future & References", 27)
    add_card_with_label(slide, Inches(0.4), Inches(1.2), Inches(4.5), Inches(2.5),
                        "\U0001f680 技术发展趋势", [
                            "模型持续变大，但更多转向效率优化 (MoE, 蒸馏)",
                            "多模态深度融合：文本+图像+语音+视频+代码",
                            "长上下文成为标配（100万Token+）",
                            "Agent能力持续增强：自主规划、执行、反思",
                            "端侧AI：手机/PC本地运行，隐私+低延迟",
                            "具身智能：大模型+机器人，迈向物理世界",
                        ])
    add_card_with_label(slide, Inches(5.2), Inches(1.2), Inches(4.5), Inches(2.5),
                        "\U0001f914 值得思考的问题", [
                            "Scaling Law会继续有效吗？还是已到瓶颈？",
                            "大模型真正「理解」语言吗？还是只是统计模式匹配？",
                            "AGI(通用人工智能)何时到来？5年？20年？",
                            "AI会成为人类的工具还是替代者？",
                            "如何确保AI的发展符合人类的长远利益？",
                            "作为未来的从业者，我们应该如何准备？",
                        ])
    # References
    add_simple_rect(slide, Inches(0.4), Inches(3.9), Inches(9.3), Inches(0.02), NAVY)
    add_textbox(slide, Inches(0.4), Inches(4.1), Inches(9.5), Inches(0.3),
                "\U0001f4da 推荐阅读与参考资源", Pt(15), NAVY, bold=True)
    add_multiline_text(slide, Inches(0.4), Inches(4.45), Inches(9.3), Inches(2.2), [
        "• [论文] Attention Is All You Need (Vaswani et al., 2017) — Transformer架构开山之作",
        "• [论文] Scaling Laws for Neural Language Models (Kaplan et al., 2020) — 规模定律的发现",
        "• [论文] Training Language Models to Follow Instructions (InstructGPT / RLHF, OpenAI 2022)",
        "• [课程] Stanford CS224N: Natural Language Processing with Deep Learning",
        "• [课程] 吴恩达 (Andrew Ng) — ChatGPT Prompt Engineering for Developers (DeepLearning.AI)",
        "• [平台] Hugging Face (huggingface.co) — 开源模型、数据集、社区",
        "• [书籍] 《大规模语言模型：从理论到实践》(复旦/清华等编著)",
    ], Pt(9), DARK_GRAY, Pt(15))
    # Q&A
    add_white_label(slide, Inches(3.0), Inches(6.5), Inches(4.0), Inches(0.55),
                    "感谢聆听  ·  欢迎提问  \U0001f64b", Pt(18))

# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════

def main():
    slides = [
        s01_cover, s02_toc, s03_paradigm_shift,
        s04_definition, s05_scaling_law, s06_emergence, s07_llm_vs_traditional,
        s08_transformer, s09_attention, s10_training, s11_rlhf, s12_tokenization,
        s13_gpt_series, s14_claude_gemini, s15_chinese_models, s16_open_source, s17_comparison,
        s18_text_gen, s19_coding, s20_industries, s21_multimodal, s22_agents,
        s23_prompt_basics, s24_prompt_advanced, s25_prompt_examples,
        s26_challenges, s27_future,
    ]
    for i, fn in enumerate(slides):
        fn()
        print(f"  [{i+1:2d}/27] {fn.__name__[4:]}")

    output_path = r"D:\First_CC\LLM_大语言模型.pptx"
    prs.save(output_path)
    print(f"\nPPT saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
